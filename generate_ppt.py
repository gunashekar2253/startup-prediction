import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    if isinstance(content, str):
        tf.text = content
    elif isinstance(content, list):
        for i, item in enumerate(content):
            if i == 0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
                if item.startswith("  -"):
                    p.level = 1
                elif item.startswith("    -"):
                    p.level = 2
    return slide

prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Start-Up Success Prediction"
subtitle.text = "Advanced Machine Learning Framework\nProject Presentation"

# Slide 2: Agenda
agenda_content = [
    "Introduction & Overview",
    "Requirement Analysis",
    "Functional & Non-Functional System Design",
    "Proposed System Features",
    "UML Use Case Diagram",
    "Data Flow Diagram (DFD)",
    "Conclusion"
]
create_slide(prs, "Agenda", agenda_content)

# Slide 3: Intro
intro_content = [
    "Problem Statement: Predicting start-up success is complex due to nonlinear relationships and diverse factors.",
    "Solution Overview: An advanced Machine Learning framework leveraging multiple algorithms to predict success viability.",
    "Goal: Provide investors and founders with data-backed actionable insights and confidence scores."
]
create_slide(prs, "Introduction & Overview", intro_content)

# Slide 4: Requirement Analysis
req_content = [
    "Requirements analysis ensures the system meets the core needs of stakeholders.",
    "  - User Requirements:",
    "    - Intuitive dashboard for entering start-up data.",
    "    - Clear, understandable predictions (white-box AI).",
    "  - Business Requirements:",
    "    - Improve investment decision accuracy.",
    "    - Provide actionable recommendations for start-ups.",
    "  - Technical Requirements:",
    "    - Real-time data integration.",
    "    - Automated model retraining (MLOps) pipeline."
]
create_slide(prs, "Requirement Analysis", req_content)

# Slide 5: System Design
design_content = [
    "System Design ensures scalability, reliability, and security of the prediction framework.",
    "  - Functional Design:",
    "    - Data Collection & Integration Engine.",
    "    - Multi-algorithm predictive modeling (RF, GBM, DNN).",
    "    - Explainable AI Module for feature importance.",
    "    - Interactive Visual Analytics Dashboard.",
    "  - Non-Functional Design:",
    "    - Performance: Real-time prediction rendering (<2 seconds).",
    "    - Scalability: Capable of handling massive diverse datasets.",
    "    - Maintainability: Automated continuous learning/retraining."
]
create_slide(prs, "Functional and Non-Functional System Design", design_content)

# Slides 6, 7, 8: Proposed Features
feat1 = [
    "1. Advanced Machine Learning Framework",
    "  - Multi-Algorithm Approach (Random Forests, Gradient Boosting, Deep Neural Networks).",
    "  - High Accuracy by capturing intricate data patterns.",
    "2. Comprehensive & Diverse Data Integration",
    "  - Rich Dataset: Financials, founder profiles, market trends, etc."
]
create_slide(prs, "Proposed System Features (Part 1/3)", feat1)

feat2 = [
    "3. Real-Time Analytics & Responsiveness",
    "  - Live Data Feeds to model current market conditions.",
    "  - Rapid reassessment as new data arrives.",
    "4. Automated Continuous Learning (MLOps)",
    "  - Dynamic Adaptation through self-improvement.",
    "  - Automated Retraining based on new trends and historical outcomes."
]
create_slide(prs, "Proposed System Features (Part 2/3)", feat2)

feat3 = [
    "5. Explainable AI (XAI) for Transparency",
    "  - White-Box predictions with direct feature importance insights.",
    "6. Interactive Visual Analytics Dashboard",
    "  - Confidence scores and data-backed Actionable Recommendations.",
    "7. Advanced Feature Engineering",
    "  - Dimensionality reduction to optimize performance."
]
create_slide(prs, "Proposed System Features (Part 3/3)", feat3)

# Slide 9: UML Diagram
blank_slide_layout = prs.slide_layouts[5] # title only
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = "System UML Diagram - Use Case Overview"

shapes = slide.shapes
left = Inches(1)
top = Inches(2)
width = Inches(1.5)
height = Inches(0.8)

actor_box = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.5), width, height)
actor_box.text = "Actor\n(Investor / Founder)"

sys_box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), top, Inches(5.5), Inches(4))
sys_box.fill.background()
sys_box.text_frame.text = ""

uc1 = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(2.2), Inches(2.5), Inches(0.8))
uc1.text = "Input Start-Up Data"

uc2 = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(3.5), Inches(2.5), Inches(0.8))
uc2.text = "View Prediction Score"

uc3 = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(4.8), Inches(2.5), Inches(0.8))
uc3.text = "View AI Explanations"

uc_db = shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(3.5), Inches(1.5), Inches(0.8))
uc_db.text = "Train Model"

tx_box = shapes.add_textbox(Inches(1), Inches(6), Inches(7.5), Inches(1))
tx_box.text = "Note: Users interact with the system to provide data and receive explainable predictions.\nThe system manages model training in the backend."

# Slide 10: DFD
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = "Data Flow Diagram (Level 0 / Context)"

shapes = slide.shapes

user = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(1.5), Inches(1))
user.text = "User / Investor"

sys = shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(2.5), Inches(2), Inches(2))
sys.text = "Start-Up\nSuccess Predictor\n(ML System)"
sys.fill.solid()
sys.fill.fore_color.rgb = RGBColor(0, 102, 204)

ds = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(0.5), Inches(2), Inches(1))
ds.text = "Market / Financial\nData Sources"

rep = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(3), Inches(1.5), Inches(1))
rep.text = "Prediction\nDecisions & Output"

tx = shapes.add_textbox(Inches(2.2), Inches(2.5), Inches(1.5), Inches(0.5))
tx.text = "Start-Up Specs >"

tx2 = shapes.add_textbox(Inches(2.2), Inches(3.5), Inches(1.5), Inches(0.5))
tx2.text = "< UI Dashboard"

tx3 = shapes.add_textbox(Inches(4.2), Inches(1.6), Inches(1.5), Inches(0.5))
tx3.text = "v Live Data"

tx4 = shapes.add_textbox(Inches(6.2), Inches(3.2), Inches(1.5), Inches(0.5))
tx4.text = "Confidence Scores >"

# Slide 11: Conclusion
conc_content = [
    "Summary:",
    "  - A comprehensive, ML-driven approach to predict start-up success.",
    "  - Integrates diverse data sources for high-accuracy predictions.",
    "  - Prioritizes transparency (Explainable AI) to build stakeholder trust.",
    "  - Employs an automated retraining loop to adapt to current market trends.",
    "Impact:",
    "  - Empowers investors to make data-backed, confident funding decisions.",
    "  - Provides actionable insights for founders to improve growth viability."
]
create_slide(prs, "Conclusion", conc_content)

ppt_path = 'c:/Users/Dell/Desktop/13/Presentation_Startup_Success.pptx'
prs.save(ppt_path)
print(f"Successfully created {ppt_path}")
